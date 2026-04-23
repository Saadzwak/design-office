"""Client-pitch PowerPoint renderer for the Justify argumentaire.

Complements `_render_client_pdf` in `justify.py`: same sourced content,
different format. The PPTX is 16:9, 6 slides in the Design Office palette,
safe to hand to a client meeting or a steering committee.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.models import VariantOutput
from app.surfaces.justify import OUT_DIR as PDF_OUT_DIR  # reuse dir layout

PPTX_OUT_DIR = Path(__file__).resolve().parent.parent / "out" / "justify_pptx"

# Organic Modern palette — ivory paper deck, forest accent, sand rules.
# Mirrors frontend/tailwind.config.ts tokens so the exported deck reads
# like a continuation of the app rather than a separate artefact.
#
# The legacy names (INK = background, BONE_TEXT = default text) are kept
# so downstream code doesn't churn — but the RGB values flip the deck
# from a dark SaaS-hero aesthetic to an editorial ivory page.
INK = RGBColor(0xFA, 0xF7, 0xF2)         # canvas — slide background (was near-black)
INK_SOFT = RGBColor(0xE8, 0xE3, 0xD8)    # hairline — faint bands / rules
BONE = RGBColor(0xFF, 0xFC, 0xF6)        # raised — card fills
BONE_TEXT = RGBColor(0x1C, 0x1F, 0x1A)   # ink — default text on ivory
TERRACOTTA = RGBColor(0x2F, 0x4A, 0x3F)  # forest — accent for eyebrows / headlines
OCHRE = RGBColor(0xA0, 0x88, 0x63)       # sand-deep — hairlines, rules
NEUTRAL_400 = RGBColor(0x7F, 0x83, 0x7D)  # ink-muted — secondary labels
NEUTRAL_300 = RGBColor(0xC6, 0xC1, 0xB4)  # mist-300 — tertiary hairlines


@dataclass(frozen=True)
class PptxBuild:
    pptx_id: str
    path: Path
    slide_count: int
    bytes: int


def render_pitch_deck(
    *,
    client_name: str,
    variant: VariantOutput,
    argumentaire_markdown: str,
    project_reference: str | None = None,
    client_logo_data_url: str | None = None,
    sketchup_iso_path: str | None = None,
    # Iter-20e additions — optional rich inputs. When present the new
    # slides render ; when absent the deck still lands the core 6.
    tagline: str | None = None,
    palette_hexes: list[str] | None = None,
    programme_markdown: str | None = None,
    other_variants: list[VariantOutput] | None = None,
    sketchup_iso_by_style: dict[str, str] | None = None,
    gallery_tile_paths: dict[str, str] | None = None,
    materials: list[dict] | None = None,
    furniture: list[dict] | None = None,
) -> PptxBuild:
    """Render the client pitch deck — 6 to 12+ slides depending on the
    optional rich inputs (iter-20e, Saad #19-#22).

    Optional extras :

    - `client_logo_data_url` : data URL from the Brief page.
    - `sketchup_iso_path` : absolute path to a PNG iso render of the
      retained variant. Embedded on the cover + on the retained-variant
      focus slide.
    - `tagline` : one-liner from the mood-board curator (used on the
      Vision slide under the client name).
    - `palette_hexes` : up to 6 hex strings for the palette strip on
      Vision + Atmosphere slides.
    - `programme_markdown` : the brief's consolidated programme. The
      "Programme" slide renders the first 6-8 H2 sections as cards.
    - `other_variants` + `sketchup_iso_by_style` : used on the "Three
      variants" slide to show all 3 macro-zoning variants in a strip.
    - `gallery_tile_paths` : `{"atmosphere": ".../a.png", "materials":
      ".../m.png", "furniture": ".../f.png", "biophilic": ".../b.png"}`
      — paths to the NanoBanana gallery tiles (iter-20d). Used on
      Atmosphere + Materials slides.
    - `materials` / `furniture` : lists from the mood-board curator
      (used on the Materials slide as a caption strip).
    """

    PPTX_OUT_DIR.mkdir(parents=True, exist_ok=True)
    signature = (
        f"pptx:{client_name}:{variant.style.value}:"
        f"{argumentaire_markdown[:400]}:"
        f"{bool(client_logo_data_url)}:{bool(sketchup_iso_path)}:"
        f"{tagline or ''}:{','.join(palette_hexes or [])}:"
        f"{len(other_variants or [])}:"
        f"{','.join(sorted((gallery_tile_paths or {}).keys()))}"
    )
    pptx_id = hashlib.sha1(signature.encode("utf-8")).hexdigest()[:16]
    target = PPTX_OUT_DIR / f"{pptx_id}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    sections = _split_argumentaire(argumentaire_markdown)
    logo_bytes = _decode_data_url(client_logo_data_url) if client_logo_data_url else None

    _build_cover_slide(
        prs,
        client_name=client_name,
        variant=variant,
        logo_bytes=logo_bytes,
        iso_path=sketchup_iso_path,
    )
    # II. Vision — tagline + palette strip (NEW iter-20e).
    _build_vision_slide(
        prs,
        client_name=client_name,
        variant=variant,
        tagline=tagline,
        palette_hexes=palette_hexes or [],
    )
    # III. Programme — condensed from markdown (NEW iter-20e).
    _build_programme_slide(prs, programme_markdown=programme_markdown or "")
    # IV. Three macro variants (NEW iter-20e).
    _build_three_variants_slide(
        prs,
        retained=variant,
        other_variants=other_variants or [],
        iso_by_style=sketchup_iso_by_style or {},
    )
    # V. Retained variant focus (NEW iter-20e).
    _build_retained_variant_slide(
        prs,
        variant=variant,
        iso_path=sketchup_iso_path,
    )
    # VI. Le pari / The bet (existing).
    _build_bet_slide(prs, client_name=client_name, section=sections.get("1", ""))
    # VII. Programme metrics (existing).
    _build_metrics_slide(prs, variant=variant)
    # VIII. Atmosphere mood board (NEW iter-20e).
    _build_atmosphere_slide(
        prs,
        tagline=tagline,
        palette_hexes=palette_hexes or [],
        atmosphere_image=(gallery_tile_paths or {}).get("atmosphere"),
        biophilic_image=(gallery_tile_paths or {}).get("biophilic"),
    )
    # IX. Materials & furniture mood board (NEW iter-20e).
    _build_materials_furniture_slide(
        prs,
        materials_image=(gallery_tile_paths or {}).get("materials"),
        furniture_image=(gallery_tile_paths or {}).get("furniture"),
        materials=materials or [],
        furniture=furniture or [],
    )
    # X. Research citations (existing).
    _build_research_slide(prs, sections=sections)
    # XI. Regulatory + constraints (existing).
    _build_regulatory_slide(prs, sections=sections)
    # XII. Next steps (existing).
    _build_next_steps_slide(
        prs,
        sections=sections,
        project_reference=project_reference,
        logo_bytes=logo_bytes,
    )

    prs.save(str(target))
    size = target.stat().st_size

    return PptxBuild(
        pptx_id=pptx_id,
        path=target,
        slide_count=len(prs.slides),
        bytes=size,
    )


def _decode_data_url(data_url: str) -> bytes | None:
    """Extract bytes from a `data:image/png;base64,…` URL."""

    import base64

    try:
        head, tail = data_url.split(",", 1)
        if "base64" not in head:
            return None
        return base64.b64decode(tail)
    except Exception:  # noqa: BLE001
        return None


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _blank_slide(prs: Presentation) -> tuple:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Paint the background ink.
    _add_background(slide, INK, prs.slide_width, prs.slide_height)
    return slide, prs.slide_width, prs.slide_height


def _add_background(slide, fill: RGBColor, w, h) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill


def _add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font: str = "Inter",
    size: int = 14,
    color: RGBColor = BONE_TEXT,
    bold: bool = False,
    italic: bool = False,
    letter_spacing: int | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_hr(slide, left: float, top: float, width: float, color: RGBColor = OCHRE) -> None:
    line = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + width), Inches(top))
    line.line.color.rgb = color
    line.line.width = Pt(1.0)


def _build_cover_slide(
    prs: Presentation,
    *,
    client_name: str,
    variant: VariantOutput,
    logo_bytes: bytes | None = None,
    iso_path: str | None = None,
) -> None:
    slide, _, _ = _blank_slide(prs)

    # Optional SketchUp iso render on the right half of the cover.
    if iso_path and Path(iso_path).exists():
        try:
            slide.shapes.add_picture(
                iso_path,
                Inches(7.0),
                Inches(0.6),
                width=Inches(5.8),
                height=Inches(6.3),
            )
        except Exception:  # noqa: BLE001
            pass

    # Optional client logo, top-left.
    if logo_bytes:
        try:
            import io as _io

            slide.shapes.add_picture(
                _io.BytesIO(logo_bytes),
                Inches(0.8),
                Inches(0.5),
                height=Inches(0.6),
            )
        except Exception:  # noqa: BLE001
            pass

    _add_text(
        slide,
        0.8,
        1.3,
        6.0,
        0.35,
        "Design Office — client argumentaire",
        font="Courier New",
        size=11,
        color=TERRACOTTA,
        letter_spacing=250,
    )
    _add_text(
        slide,
        0.8,
        1.9,
        6.0,
        2.6,
        f"{client_name}\n« {variant.title} »",
        font="Calibri",
        size=48,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 4.6, 4.0)
    _add_text(
        slide,
        0.8,
        4.8,
        6.0,
        2.0,
        f"Parti: {variant.style.value.replace('_', ' ').title()}\n"
        f"Desks: {variant.metrics.workstation_count}\n"
        f"Flex ratio: {variant.metrics.flex_ratio_applied:.2f}\n"
        f"Total programmed: ≈ {round(variant.metrics.total_programmed_m2)} m²",
        font="Calibri",
        size=18,
        color=BONE_TEXT,
    )
    _add_text(
        slide,
        0.8,
        6.9,
        12,
        0.4,
        f"{datetime.now(tz=timezone.utc).date().isoformat()} · Built with Opus 4.7 · MIT License",
        font="Courier New",
        size=10,
        color=NEUTRAL_400,
    )


def _build_bet_slide(prs: Presentation, *, client_name: str, section: str) -> None:
    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "01 · The bet")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        f"Why this variant for {client_name}",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)
    body = _condense(section, max_chars=1100)
    _add_text(
        slide, 0.8, 2.6, 12, 4.4, body, font="Calibri", size=18, color=BONE_TEXT
    )


def _build_metrics_slide(prs: Presentation, *, variant: VariantOutput) -> None:
    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "02 · Retained programme")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "Chiffres clés de la variante",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    metrics = [
        ("Postes individuels", str(variant.metrics.workstation_count)),
        ("Réunions", str(variant.metrics.meeting_room_count)),
        ("Phone booths", str(variant.metrics.phone_booth_count)),
        ("Flex ratio", f"{variant.metrics.flex_ratio_applied:.2f}"),
        ("Collab (m²)", f"{round(variant.metrics.collab_surface_m2)}"),
        ("Total programmed (m²)", f"{round(variant.metrics.total_programmed_m2)}"),
    ]

    x0, y0 = 0.8, 2.8
    col_w, col_h = 4.1, 1.7
    gutter = 0.2
    for i, (label, value) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left = x0 + col * (col_w + gutter)
        top = y0 + row * (col_h + gutter)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(col_w),
            Inches(col_h),
        )
        card.adjustments[0] = 0.06
        card.line.color.rgb = NEUTRAL_400
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = INK_SOFT

        _add_text(
            slide,
            left + 0.3,
            top + 0.25,
            col_w - 0.6,
            0.35,
            label.upper(),
            font="Courier New",
            size=10,
            color=NEUTRAL_300,
        )
        _add_text(
            slide,
            left + 0.3,
            top + 0.65,
            col_w - 0.6,
            0.9,
            value,
            font="Calibri",
            size=38,
            color=BONE_TEXT,
            bold=True,
        )


def _build_research_slide(prs: Presentation, *, sections: dict[str, str]) -> None:
    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "03 · What the research says")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "Acoustique · Biophilie · Flex",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    body = _condense(sections.get("2", ""), max_chars=1400)
    _add_text(
        slide, 0.8, 2.6, 12, 4.4, body, font="Calibri", size=15, color=BONE_TEXT
    )
    _add_text(
        slide,
        0.8,
        6.7,
        12,
        0.5,
        "Sources : Browning 2014 · Nieuwenhuis 2014 · Ulrich 1984 · NF S 31-080 · Leesman 2024",
        font="Courier New",
        size=9,
        color=NEUTRAL_300,
    )


def _build_regulatory_slide(prs: Presentation, *, sections: dict[str, str]) -> None:
    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "04 · What the regulation says")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "ERP type W · PMR · Code du travail",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    body = _condense(sections.get("3", ""), max_chars=1500)
    _add_text(
        slide, 0.8, 2.6, 12, 4.4, body, font="Calibri", size=15, color=BONE_TEXT
    )
    _add_text(
        slide,
        0.8,
        6.7,
        12,
        0.5,
        "Sources : Arrêté 20 avril 2017 · Règlement sécurité ERP type W · R. 4222 / R. 4223 · EN 12464-1",
        font="Courier New",
        size=9,
        color=NEUTRAL_300,
    )


def _build_next_steps_slide(
    prs: Presentation,
    *,
    sections: dict[str, str],
    project_reference: str | None,
    logo_bytes: bytes | None = None,
) -> None:
    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "05 · Next steps & KPIs")
    # Client logo, bottom-right footer.
    if logo_bytes:
        try:
            import io as _io

            slide.shapes.add_picture(
                _io.BytesIO(logo_bytes),
                Inches(12.0),
                Inches(6.7),
                height=Inches(0.5),
            )
        except Exception:  # noqa: BLE001
            pass
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "What it takes to start",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    kpis = _condense(sections.get("5", ""), max_chars=600)
    steps = _condense(sections.get("6", ""), max_chars=700)

    _add_text(
        slide,
        0.8,
        2.7,
        5.9,
        0.35,
        "EXPECTED RESULTS 6–12 MONTHS",
        font="Courier New",
        size=10,
        color=OCHRE,
    )
    _add_text(slide, 0.8, 3.1, 5.9, 3.4, kpis, font="Calibri", size=14, color=BONE_TEXT)

    _add_text(
        slide,
        7.0,
        2.7,
        5.9,
        0.35,
        "NEXT STEPS",
        font="Courier New",
        size=10,
        color=OCHRE,
    )
    _add_text(slide, 7.0, 3.1, 5.9, 3.4, steps, font="Calibri", size=14, color=BONE_TEXT)

    footer = f"Project: {project_reference or 'DO-CAT-B'} · Built with Opus 4.7"
    _add_text(
        slide,
        0.8,
        6.85,
        12,
        0.4,
        footer,
        font="Courier New",
        size=9,
        color=NEUTRAL_300,
    )


def _eyebrow(slide, text: str) -> None:
    _add_text(
        slide,
        0.8,
        0.6,
        12,
        0.35,
        text.upper(),
        font="Courier New",
        size=10,
        color=TERRACOTTA,
    )


# ---------------------------------------------------------------------------
# Iter-20e slide builders — magazine / editorial direction
# ---------------------------------------------------------------------------


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Parse `#RRGGBB` (or `RRGGBB`) → RGBColor. Defaults to ink on failure."""

    try:
        h = hex_str.lstrip("#").strip()
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        if len(h) != 6:
            return BONE_TEXT
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:  # noqa: BLE001
        return BONE_TEXT


def _add_palette_strip(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    hexes: list[str],
) -> None:
    """Draw a horizontal strip of up to 6 color chips (with hex caption)."""

    if not hexes:
        return
    chips = hexes[:6]
    chip_w = width / len(chips)
    for i, hx in enumerate(chips):
        chip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + i * chip_w),
            Inches(top),
            Inches(chip_w - 0.05),
            Inches(height),
        )
        chip.line.fill.background()
        chip.fill.solid()
        chip.fill.fore_color.rgb = _hex_to_rgb(hx)
        _add_text(
            slide,
            left + i * chip_w,
            top + height + 0.08,
            chip_w - 0.05,
            0.3,
            hx.upper(),
            font="Courier New",
            size=8,
            color=NEUTRAL_400,
        )


def _safe_add_picture(
    slide,
    image_path: str | None,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    placeholder_label: str = "",
) -> None:
    """Add a picture if the file exists; otherwise draw an ivory placeholder.

    Keeps the deck visually honest: empty frame with a label instead of
    a missing image.
    """

    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height),
            )
            return
        except Exception:  # noqa: BLE001
            pass
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame.line.color.rgb = NEUTRAL_300
    frame.line.width = Pt(0.5)
    frame.fill.solid()
    frame.fill.fore_color.rgb = INK_SOFT
    if placeholder_label:
        _add_text(
            slide,
            left + 0.25,
            top + height / 2 - 0.2,
            width - 0.5,
            0.4,
            placeholder_label.upper(),
            font="Courier New",
            size=9,
            color=NEUTRAL_400,
        )


def _build_vision_slide(
    prs: Presentation,
    *,
    client_name: str,
    variant: VariantOutput,
    tagline: str | None,
    palette_hexes: list[str],
) -> None:
    """Editorial opener: tagline + client name + palette strip.

    Title-page of the narrative. Very little text, lots of air. Feels
    like a Kinfolk cover. No image needed — the palette carries the
    atmosphere."""

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "Vision")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "The intent",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    headline = (tagline or "").strip()
    if not headline:
        headline = f"Designing {client_name}'s {variant.title.lower()}."
    # Big editorial pull quote.
    _add_text(
        slide,
        0.8,
        2.7,
        11.8,
        2.4,
        f"« {headline} »",
        font="Calibri",
        size=34,
        color=BONE_TEXT,
        italic=True,
    )

    signature = f"{client_name} · {variant.title}"
    _add_text(
        slide,
        0.8,
        5.2,
        11.8,
        0.4,
        signature,
        font="Courier New",
        size=11,
        color=OCHRE,
    )

    if palette_hexes:
        _add_text(
            slide,
            0.8,
            5.85,
            11.8,
            0.3,
            "PALETTE",
            font="Courier New",
            size=9,
            color=NEUTRAL_400,
        )
        _add_palette_strip(
            slide,
            left=0.8,
            top=6.15,
            width=11.8,
            height=0.55,
            hexes=palette_hexes,
        )


def _build_programme_slide(
    prs: Presentation,
    *,
    programme_markdown: str,
) -> None:
    """Condensed programme — first 6 H2 sections rendered as caption cards.

    Reads the consolidator's markdown, extracts `## N. Title` headings
    plus their opening sentence, lays them out in a 3×2 grid. Gives
    the client a single-glance read of what was programmed.
    """

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "Programme")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "What we programmed",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    entries = _extract_programme_entries(programme_markdown, limit=6)
    if not entries:
        _add_text(
            slide,
            0.8,
            2.7,
            12,
            1.0,
            "Programme à consolider à partir du brief.",
            font="Calibri",
            size=14,
            color=NEUTRAL_400,
            italic=True,
        )
        return

    x0, y0 = 0.8, 2.7
    col_w, col_h = 4.0, 1.9
    gutter = 0.15
    for i, (title, body) in enumerate(entries):
        col = i % 3
        row = i // 3
        left = x0 + col * (col_w + gutter)
        top = y0 + row * (col_h + gutter)
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(col_w),
            Inches(col_h),
        )
        card.line.color.rgb = NEUTRAL_300
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = BONE
        _add_text(
            slide,
            left + 0.25,
            top + 0.2,
            col_w - 0.5,
            0.35,
            title.upper(),
            font="Courier New",
            size=9,
            color=TERRACOTTA,
        )
        _add_text(
            slide,
            left + 0.25,
            top + 0.6,
            col_w - 0.5,
            col_h - 0.7,
            body,
            font="Calibri",
            size=11,
            color=BONE_TEXT,
        )


def _build_three_variants_slide(
    prs: Presentation,
    *,
    retained: VariantOutput,
    other_variants: list[VariantOutput],
    iso_by_style: dict[str, str],
) -> None:
    """Strip of 3 variants — retained in bold, others dimmer.

    One column per variant, 3D iso on top, title below, metrics below
    that, a discreet "retained" / "explored" marker. Lets the client
    see the full decision set in one glance.
    """

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "The three macro-zonings")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "Three hypotheses, one bet",
        font="Calibri",
        size=32,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.15, 4.0)

    # Order: retained first, then the two others.
    variants: list[VariantOutput] = [retained]
    for v in other_variants:
        if v.style != retained.style and len(variants) < 3:
            variants.append(v)

    col_w = 4.0
    gutter = 0.2
    x0 = 0.8
    y_img = 2.5
    img_h = 2.6
    for i, v in enumerate(variants[:3]):
        left = x0 + i * (col_w + gutter)
        iso_path = iso_by_style.get(v.style.value) if iso_by_style else None
        _safe_add_picture(
            slide,
            iso_path,
            left=left,
            top=y_img,
            width=col_w,
            height=img_h,
            placeholder_label=f"{v.style.value.replace('_', ' ')} · iso",
        )
        is_retained = v.style == retained.style
        label_color = TERRACOTTA if is_retained else NEUTRAL_400
        label_txt = "RETAINED" if is_retained else "EXPLORED"
        _add_text(
            slide,
            left,
            y_img + img_h + 0.1,
            col_w,
            0.3,
            label_txt,
            font="Courier New",
            size=9,
            color=label_color,
        )
        _add_text(
            slide,
            left,
            y_img + img_h + 0.45,
            col_w,
            0.55,
            v.title,
            font="Calibri",
            size=18,
            color=BONE_TEXT,
            bold=is_retained,
        )
        metrics = (
            f"{v.metrics.workstation_count} desks · "
            f"flex {v.metrics.flex_ratio_applied:.2f} · "
            f"{round(v.metrics.total_programmed_m2)} m²"
        )
        _add_text(
            slide,
            left,
            y_img + img_h + 1.05,
            col_w,
            0.5,
            metrics,
            font="Courier New",
            size=10,
            color=NEUTRAL_400,
        )


def _build_retained_variant_slide(
    prs: Presentation,
    *,
    variant: VariantOutput,
    iso_path: str | None,
) -> None:
    """Full-bleed focus on the chosen variant — big 3D iso + short pitch.

    Image dominates the left two-thirds, copy column on the right.
    Two key metrics as an overlay strip at the bottom of the image.
    """

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "The retained variant")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        variant.title,
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    _safe_add_picture(
        slide,
        iso_path,
        left=0.8,
        top=2.5,
        width=8.0,
        height=4.5,
        placeholder_label="sketchup iso render",
    )

    # Copy column (right).
    right_x = 9.1
    right_w = 3.6
    _add_text(
        slide,
        right_x,
        2.5,
        right_w,
        0.3,
        "PARTI",
        font="Courier New",
        size=9,
        color=TERRACOTTA,
    )
    _add_text(
        slide,
        right_x,
        2.85,
        right_w,
        1.1,
        variant.style.value.replace("_", " ").title(),
        font="Calibri",
        size=20,
        color=BONE_TEXT,
        bold=True,
    )

    _add_text(
        slide,
        right_x,
        4.0,
        right_w,
        0.3,
        "FLEX RATIO",
        font="Courier New",
        size=9,
        color=TERRACOTTA,
    )
    _add_text(
        slide,
        right_x,
        4.35,
        right_w,
        0.6,
        f"{variant.metrics.flex_ratio_applied:.2f}",
        font="Calibri",
        size=26,
        color=BONE_TEXT,
        bold=True,
    )

    _add_text(
        slide,
        right_x,
        5.2,
        right_w,
        0.3,
        "TOTAL PROGRAMMED",
        font="Courier New",
        size=9,
        color=TERRACOTTA,
    )
    _add_text(
        slide,
        right_x,
        5.55,
        right_w,
        0.6,
        f"{round(variant.metrics.total_programmed_m2)} m²",
        font="Calibri",
        size=26,
        color=BONE_TEXT,
        bold=True,
    )


def _build_atmosphere_slide(
    prs: Presentation,
    *,
    tagline: str | None,
    palette_hexes: list[str],
    atmosphere_image: str | None,
    biophilic_image: str | None,
) -> None:
    """Mood slide — two large tiles from the NanoBanana gallery.

    Left tile: atmosphere. Right tile: biophilic. Palette band below.
    Tagline as caption. This is the slide that sells the *feel* of
    the project before the numbers.
    """

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "Atmosphere")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "How it feels",
        font="Calibri",
        size=36,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.2, 4.0)

    _safe_add_picture(
        slide,
        atmosphere_image,
        left=0.8,
        top=2.5,
        width=6.0,
        height=3.8,
        placeholder_label="atmosphere tile",
    )
    _safe_add_picture(
        slide,
        biophilic_image,
        left=7.0,
        top=2.5,
        width=5.6,
        height=3.8,
        placeholder_label="biophilic tile",
    )

    if tagline:
        _add_text(
            slide,
            0.8,
            6.4,
            11.8,
            0.45,
            f"« {tagline.strip()} »",
            font="Calibri",
            size=16,
            color=BONE_TEXT,
            italic=True,
        )

    if palette_hexes:
        _add_palette_strip(
            slide,
            left=0.8,
            top=6.95,
            width=11.8,
            height=0.3,
            hexes=palette_hexes,
        )


def _build_materials_furniture_slide(
    prs: Presentation,
    *,
    materials_image: str | None,
    furniture_image: str | None,
    materials: list[dict],
    furniture: list[dict],
) -> None:
    """Split mood — two gallery tiles + concise caption strips.

    Captions pull `name · brand` (or `material`) from the curator's
    selection and render them as bullet lines. No prices, no SKUs —
    the client doesn't care, the architect does."""

    slide, _, _ = _blank_slide(prs)
    _eyebrow(slide, "Materials & furniture")
    _add_text(
        slide,
        0.8,
        1.15,
        12,
        1.0,
        "The palette, made real",
        font="Calibri",
        size=32,
        color=BONE_TEXT,
        bold=True,
    )
    _add_hr(slide, 0.8, 2.15, 4.0)

    # Left half — Materials.
    _add_text(
        slide,
        0.8,
        2.45,
        6.0,
        0.3,
        "MATERIALS",
        font="Courier New",
        size=9,
        color=TERRACOTTA,
    )
    _safe_add_picture(
        slide,
        materials_image,
        left=0.8,
        top=2.8,
        width=6.0,
        height=3.0,
        placeholder_label="materials tile",
    )
    _add_text(
        slide,
        0.8,
        5.95,
        6.0,
        1.3,
        _caption_from_list(materials, keys=("material", "name"), limit=5),
        font="Calibri",
        size=11,
        color=BONE_TEXT,
    )

    # Right half — Furniture.
    _add_text(
        slide,
        7.0,
        2.45,
        5.6,
        0.3,
        "FURNITURE",
        font="Courier New",
        size=9,
        color=TERRACOTTA,
    )
    _safe_add_picture(
        slide,
        furniture_image,
        left=7.0,
        top=2.8,
        width=5.6,
        height=3.0,
        placeholder_label="furniture tile",
    )
    _add_text(
        slide,
        7.0,
        5.95,
        5.6,
        1.3,
        _caption_from_list(furniture, keys=("name", "brand", "type"), limit=5),
        font="Calibri",
        size=11,
        color=BONE_TEXT,
    )


def _extract_programme_entries(
    md: str, *, limit: int = 6
) -> list[tuple[str, str]]:
    """Find `## Title` headings and return `(title, first-paragraph)` pairs.

    Skips markdown table rows and empty lines when picking the snippet.
    """

    if not md:
        return []
    lines = md.splitlines()
    entries: list[tuple[str, str]] = []
    current_title: str | None = None
    current_body: list[str] = []

    def _flush() -> None:
        if current_title is None:
            return
        # First non-table, non-empty, non-heading paragraph.
        snippet_parts: list[str] = []
        for raw in current_body:
            s = raw.strip()
            if not s:
                if snippet_parts:
                    break
                continue
            if s.startswith("|") or s.startswith("#"):
                continue
            # Strip markdown emphasis / links for a clean caption.
            s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
            s = re.sub(r"(?<!\*)\*([^*\n]+?)\*(?!\*)", r"\1", s)
            s = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", s)
            s = re.sub(r"`([^`]+)`", r"\1", s)
            snippet_parts.append(s)
            if sum(len(p) for p in snippet_parts) >= 180:
                break
        snippet = " ".join(snippet_parts).strip()
        if len(snippet) > 220:
            snippet = snippet[:219].rstrip() + "…"
        entries.append((current_title, snippet or ""))

    for ln in lines:
        m = re.match(r"^##\s+(?:\d+\.\s*)?(.+?)\s*$", ln)
        if m:
            _flush()
            if len(entries) >= limit:
                current_title = None
                current_body = []
                return entries
            current_title = m.group(1).strip()
            current_body = []
            continue
        if current_title is not None:
            current_body.append(ln)
    _flush()
    return entries[:limit]


def _caption_from_list(
    items: list[dict], *, keys: tuple[str, ...], limit: int = 5
) -> str:
    """Format `[{name, brand, ...}]` as a compact bullet block.

    Pulls the first populated key from `keys` for the primary label,
    and the next one as a sidekick when it exists (`Name · Brand`)."""

    if not items:
        return ""
    out: list[str] = []
    for it in items[:limit]:
        if not isinstance(it, dict):
            continue
        primary = ""
        secondary = ""
        for k in keys:
            val = str(it.get(k, "")).strip()
            if val and not primary:
                primary = val
            elif val and not secondary and val != primary:
                secondary = val
        if not primary:
            continue
        line = f"•  {primary}"
        if secondary:
            line += f"  ·  {secondary}"
        out.append(line)
    return "\n".join(out)


# ---------------------------------------------------------------------------
# Markdown helpers
# ---------------------------------------------------------------------------


_SECTION_RE = re.compile(r"^##\s+(\d+)\.\s*(.+?)\s*$", flags=re.MULTILINE)


def _split_argumentaire(md: str) -> dict[str, str]:
    """Split the consolidator Markdown by its numbered `## N. Title` headings.

    Returns {"1": "...body...", "2": "...", ...}. Each body is the raw
    Markdown between that heading and the next one (or EOF).
    """

    out: dict[str, str] = {}
    matches = list(_SECTION_RE.finditer(md))
    for i, m in enumerate(matches):
        idx = m.group(1)
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md)
        out[idx] = md[start:end].strip()
    return out


def _condense(raw: str, *, max_chars: int) -> str:
    """Turn Markdown subsection body into slide-friendly plain text.

    Strips `### X` subheadings (keeping a separator line), drops bold/italic
    markers, preserves `- ` bullet markers as `•`, truncates to `max_chars`
    with an ellipsis.
    """

    if not raw:
        return ""
    text = raw
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*([^*\n]+?)\*(?!\*)", r"\1", text)
    text = re.sub(r"^###\s+(.+?)\s*$", r"— \1", text, flags=re.MULTILINE)
    text = re.sub(r"^-\s+", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()
    if len(text) > max_chars:
        text = text[: max_chars - 1].rstrip() + "…"
    return text


def pptx_path_for(pptx_id: str) -> Path | None:
    candidate = PPTX_OUT_DIR / f"{pptx_id}.pptx"
    return candidate if candidate.exists() else None


def pdf_out_dir() -> Path:
    """Helper to keep `main.py` tidy — the Justify PDF directory lives in
    `justify.py` as the canonical source of truth.
    """

    return PDF_OUT_DIR
